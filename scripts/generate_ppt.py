"""Generate PPT Brief for TA GraphRAG-Kubernetes."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette ──────────────────────────────────────────────────────────────
C_DARK_BLUE   = RGBColor(0x1A, 0x29, 0x4E)   # navy — slide background / header
C_ACCENT_TEAL = RGBColor(0x00, 0x9B, 0x8D)   # teal — accent / highlight
C_LIGHT_BG    = RGBColor(0xF5, 0xF7, 0xFA)   # near-white — body bg
C_TEXT_DARK   = RGBColor(0x1A, 0x29, 0x4E)   # navy — body text
C_TEXT_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TABLE_HDR   = RGBColor(0x00, 0x9B, 0x8D)   # teal header row
C_TABLE_ALT   = RGBColor(0xE8, 0xF6, 0xF5)   # light teal alt row
C_GOLD        = RGBColor(0xFF, 0xC8, 0x00)   # gold — emphasis


# ── Helpers ────────────────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txbox(slide, text, left, top, width, height,
          font_size=18, bold=False, color=C_TEXT_DARK,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tb.text_frame.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_header_bar(slide, title: str, subtitle: str = ""):
    """Dark blue header bar at top of content slides."""
    add_rect(slide, 0, 0, 13.33, 1.5, C_DARK_BLUE)
    txbox(slide, title, 0.3, 0.15, 12.5, 0.8,
          font_size=28, bold=True, color=C_TEXT_WHITE)
    if subtitle:
        txbox(slide, subtitle, 0.3, 0.85, 12.5, 0.5,
              font_size=14, color=C_ACCENT_TEAL)
    # accent line
    add_rect(slide, 0, 1.5, 13.33, 0.05, C_ACCENT_TEAL)


def bullet_block(slide, items, left, top, width, height,
                 font_size=16, color=C_TEXT_DARK, indent_char="•  "):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = indent_char + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return tb


def add_table(slide, headers, rows, left, top, width,
              col_widths=None, font_size=13):
    """Add a styled table."""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 header
    row_h = Inches(0.38)
    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), row_h * n_rows
    ).table

    # Column widths
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(width * cw / total)

    def _cell(r, c, text, bg, fg=C_TEXT_WHITE, bold=False, sz=font_size):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = fg

    for c, h in enumerate(headers):
        _cell(0, c, h, C_TABLE_HDR, C_TEXT_WHITE, bold=True)

    for r, row in enumerate(rows):
        bg = C_TABLE_ALT if r % 2 == 0 else C_TEXT_WHITE
        for c, val in enumerate(row):
            _cell(r + 1, c, val, bg, C_TEXT_DARK, sz=font_size - 1)

    return tbl


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, C_DARK_BLUE)

    # decorative teal band
    add_rect(slide, 0, 2.8, 13.33, 0.15, C_ACCENT_TEAL)
    add_rect(slide, 0, 5.8, 13.33, 0.15, C_ACCENT_TEAL)
    add_rect(slide, 0, 5.95, 13.33, 1.55, RGBColor(0x0D, 0x1A, 0x35))

    # main title
    txbox(slide,
          "Implementasi Graph Retrieval Augmented Generation\nuntuk Meningkatkan Presisi Retrieval dan Validitas Sintaksis\npada Konfigurasi Kubernetes",
          0.6, 1.0, 12.0, 1.7,
          font_size=26, bold=True, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)

    # subtitle bar
    txbox(slide, "TUGAS AKHIR  ·  SISTEM DAN TEKNOLOGI INFORMASI  ·  STEI ITB",
          0.6, 2.85, 12.0, 0.5,
          font_size=13, color=C_ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # author block
    txbox(slide, "Jihan Aurelia  ·  18222001",
          0.6, 3.4, 12.0, 0.5,
          font_size=18, bold=True, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)
    txbox(slide, "Pembimbing: Dr. Ir. Dimitri Mahayana, M.Eng.  ·  2025",
          0.6, 3.9, 12.0, 0.4,
          font_size=13, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

    # tag line
    txbox(slide, "GraphRAG  ·  Neo4j  ·  LangGraph  ·  GPT-4o-mini  ·  Groq LLaMA",
          0.6, 6.2, 12.0, 0.5,
          font_size=13, color=C_ACCENT_TEAL, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_LIGHT_BG)
    add_header_bar(slide,
                   "Latar Belakang & Masalah",
                   "Mengapa Vector RAG tidak cukup untuk Kubernetes?")

    # Left column — pain points
    txbox(slide, "Tantangan Kubernetes", 0.4, 1.7, 5.5, 0.4,
          font_size=15, bold=True, color=C_DARK_BLUE)
    bullet_block(slide, [
        "Skema hierarkis kompleks: Pod ← Deployment ← HPA",
        "Relasi referensial: Volume, PVC, StorageClass, Secret, ConfigMap",
        "Vector RAG hanya mencocokkan teks, tidak memahami relasi antar-objek",
        "YAML yang dihasilkan sering tidak valid secara sintaksis",
    ], 0.4, 2.1, 5.8, 2.5, font_size=14)

    # Right column — baseline numbers
    add_rect(slide, 6.5, 1.65, 6.4, 3.2, C_DARK_BLUE)
    txbox(slide, "Hasil Awal Vector RAG", 6.7, 1.75, 6.0, 0.45,
          font_size=15, bold=True, color=C_ACCENT_TEAL)

    stats = [
        ("Exact Match Accuracy", "63,4%"),
        ("Context Precision",    "69,8%"),
        ("YAML Syntactic Validity", "< 100%"),
    ]
    y = 2.3
    for label, val in stats:
        txbox(slide, label, 6.7, y, 4.0, 0.4, font_size=13, color=C_TEXT_WHITE)
        txbox(slide, val,   10.5, y, 2.2, 0.4, font_size=20, bold=True,
              color=C_GOLD, align=PP_ALIGN.RIGHT)
        y += 0.7

    txbox(slide, "→ Dibutuhkan pendekatan berbasis graph untuk memahami\n    relasi dan dependensi antar-resource Kubernetes",
          0.4, 4.7, 12.5, 0.7, font_size=14, italic=True, color=C_DARK_BLUE)


def slide_rq(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "Rumusan Masalah & Tujuan", "Tiga Pertanyaan Penelitian (RQ1–RQ3)")

    headers = ["RQ", "Pertanyaan Penelitian", "Target / Output"]
    rows = [
        ["RQ1", "Bagaimana membangun knowledge graph dari spesifikasi Kubernetes yang merepresentasikan struktur dan relasi antar-objek?",
         "KG dari swagger.json v1.30:\n725 node, 18 jenis edge"],
        ["RQ2", "Bagaimana merancang mekanisme GraphRAG berbasis KG untuk meningkatkan presisi retrieval?",
         "Cascading retrieval + multi-hop graph traversal adaptif per intent"],
        ["RQ3", "Bagaimana perbandingan GraphRAG vs. Vector RAG vs. Vanilla LLM?",
         "Evaluasi 3 dimensi (AnsQ/RetQ/ReaQ) pada 97 fixture"],
    ]
    add_table(slide, headers, rows,
              left=0.4, top=1.7, width=12.5,
              col_widths=[1, 5, 4], font_size=13)

    txbox(slide, "Batasan: swagger.json v1.30 · definisi schema saja · maks. 3 hop · tanpa fine-tuning · output YAML (tidak diuji di klaster nyata)",
          0.4, 5.7, 12.5, 0.55, font_size=12, italic=True,
          color=RGBColor(0x66, 0x77, 0x88))


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "Arsitektur Sistem", "Pipeline LangGraph 5-Node dengan Dual-LLM")

    # Pipeline nodes row
    nodes = [
        ("memory", "SQLite\nRiwayat"),
        ("thinker", "GPT-4o-mini\nIntent Extraction"),
        ("retriever", "Neo4j\nGraph Traversal"),
        ("speaker", "Groq LLaMA\nGenerasi Respons"),
        ("saver", "SQLite\nSimpan Sesi"),
    ]
    colors = [C_DARK_BLUE, C_ACCENT_TEAL, RGBColor(0xE6, 0x7E, 0x22),
              C_ACCENT_TEAL, C_DARK_BLUE]

    box_w, box_h = 2.1, 1.0
    gap = 0.32
    start_x = 0.3
    y_box = 1.8

    for i, ((name, desc), col) in enumerate(zip(nodes, colors)):
        x = start_x + i * (box_w + gap)
        add_rect(slide, x, y_box, box_w, box_h, col)
        txbox(slide, name, x, y_box + 0.05, box_w, 0.35,
              font_size=13, bold=True, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)
        txbox(slide, desc, x, y_box + 0.4, box_w, 0.55,
              font_size=11, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            arrow_x = x + box_w + 0.04
            txbox(slide, "→", arrow_x, y_box + 0.3, 0.25, 0.4,
                  font_size=18, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

    # User ↕ Streamlit
    txbox(slide, "Pengguna  ↕  Streamlit UI", 0.3, 1.35, 12.5, 0.4,
          font_size=13, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

    # Component table
    headers = ["Komponen", "Teknologi", "Peran"]
    rows = [
        ["Frontend",       "Streamlit",          "Antarmuka percakapan; UUID unik per sesi"],
        ["Orchestrator",   "LangGraph",           "DAG 5-node; state AgentState"],
        ["Knowledge Store","Neo4j",               "Graph + vector index + embedding cosine"],
        ["Session Store",  "SQLite",              "Riwayat percakapan antar-giliran"],
        ["Thinker LLM",    "GPT-4o-mini (T=0,0)", "Ekstrak intent → JSON terstruktur"],
        ["Speaker LLM",    "Groq LLaMA-3.1-8b (T=0,1)", "Hasilkan respons naratif / YAML"],
    ]
    add_table(slide, headers, rows,
              left=0.4, top=3.1, width=12.5,
              col_widths=[2.5, 3.5, 6], font_size=12)


def slide_kg(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "Pembangunan Knowledge Graph",
                   "Dataset: swagger.json Kubernetes v1.30 → Neo4j")

    # Left — stats
    txbox(slide, "Statistik KG", 0.4, 1.7, 5.5, 0.4,
          font_size=15, bold=True, color=C_DARK_BLUE)
    add_table(slide,
              ["Atribut", "Nilai"],
              [
                  ["Sumber", "swagger.json v1.30 (3,67 MB)"],
                  ["Definisi awal", "730 (5 noise dikecualikan)"],
                  ["Node Definition", "725"],
                  ["Jenis Edge", "18 (7 kategori)"],
                  ["Model Embedding", "text-embedding-3-small, 1.536 dim"],
                  ["Vector Index", "Neo4j native cosine similarity"],
              ],
              left=0.4, top=2.1, width=6.0,
              col_widths=[3, 3], font_size=12)

    # Right — edge taxonomy
    txbox(slide, "18 Jenis Edge — 7 Kategori", 6.7, 1.7, 6.2, 0.4,
          font_size=15, bold=True, color=C_DARK_BLUE)
    add_table(slide,
              ["Kategori", "n", "Edge Type"],
              [
                  ["Struktural",   "4", "HAS_PROPERTY, EXTENDS, ONE_OF, ANY_OF"],
                  ["Workload",     "3", "CONTAINS_POD_TEMPLATE, CONTAINS_JOB_TEMPLATE, HAS_CONTAINER"],
                  ["Penyimpanan",  "3", "CLAIMS_VOLUME, MOUNTS_VOLUME, USES_STORAGE_CLASS"],
                  ["Konfigurasi",  "2", "LOADS_CONFIGMAP, USES_SECRET"],
                  ["Jaringan",     "2", "SELECTS_POD, ROUTES_TO_SERVICE"],
                  ["RBAC",         "3", "BINDS_ROLE, BINDS_SERVICE_ACCOUNT, USES_SERVICE_ACCOUNT"],
                  ["Autoscaling",  "1", "SCALES_RESOURCE"],
              ],
              left=6.7, top=2.1, width=6.2,
              col_widths=[2.5, 0.8, 6], font_size=11)

    # Pipeline ingestion steps
    txbox(slide, "Pipeline Ingestion (5 Pass):", 0.4, 5.35, 12.5, 0.35,
          font_size=13, bold=True, color=C_DARK_BLUE)
    passes = ["Pass 1: 725 Node Creation",
              "Pass 1.5: Embedding Generation",
              "Pass 2: Structural Edges (HAS_PROPERTY)",
              "Pass 2.5: Inheritance Edges (EXTENDS / ONE_OF / ANY_OF)",
              "Pass 3: 14 Semantic Edges → Neo4j KG + Vector Index"]
    colors_p = [C_DARK_BLUE, C_DARK_BLUE, C_ACCENT_TEAL, C_ACCENT_TEAL, RGBColor(0xE6, 0x7E, 0x22)]
    bw = 2.35
    for i, (p, c) in enumerate(zip(passes, colors_p)):
        x = 0.35 + i * (bw + 0.1)
        add_rect(slide, x, 5.75, bw, 0.75, c)
        txbox(slide, p, x + 0.05, 5.78, bw - 0.1, 0.7,
              font_size=10, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            txbox(slide, "→", x + bw + 0.01, 5.95, 0.12, 0.3,
                  font_size=14, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)


def slide_retrieval(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "Mekanisme Retrieval Bertingkat",
                   "Cascading Retrieval: Exact Match → Vector → Graph Traversal")

    # Flow diagram (boxes)
    flow = [
        (C_DARK_BLUE,   "Intent JSON\n(dari Thinker)",           0.4,  1.8, 2.0, 0.9),
        (C_ACCENT_TEAL, "Phase 1\nExact Name Match",             3.0,  1.8, 2.2, 0.9),
        (C_ACCENT_TEAL, "Phase 2 (fallback)\nVector Similarity", 3.0,  3.2, 2.2, 0.9),
        (RGBColor(0xE6,0x7E,0x22), "Graph Traversal\n(SCHEMA_DEPS_QUERY)", 6.0, 2.4, 2.4, 0.9),
        (C_DARK_BLUE,   "Reasoning Path\n('A -[REL]→ B')",       9.0,  2.4, 2.2, 0.9),
        (C_DARK_BLUE,   "graph_context\n(maks. 12.000 char)",   11.5,  2.4, 1.7, 0.9),
    ]
    for col, txt, x, y, w, h in flow:
        add_rect(slide, x, y, w, h, col)
        txbox(slide, txt, x + 0.05, y + 0.1, w - 0.1, h - 0.1,
              font_size=12, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)

    # arrows
    arrows = [(2.4, 2.25, "→"), (5.2, 2.25, "ditemukan →"), (5.2, 3.65, "fallback →"),
              (8.4, 2.88, "→"), (11.2, 2.88, "→")]
    for x, y, txt in arrows:
        txbox(slide, txt, x, y, 0.9, 0.3, font_size=10,
              color=C_TEXT_DARK, align=PP_ALIGN.CENTER)

    txbox(slide, "tidak ditemukan ↓", 4.05, 2.75, 1.5, 0.35,
          font_size=10, color=C_TEXT_DARK, align=PP_ALIGN.CENTER)

    # Depth table
    txbox(slide, "Kedalaman Traversal per Intent Type", 0.4, 4.55, 12.5, 0.4,
          font_size=14, bold=True, color=C_DARK_BLUE)
    add_table(slide,
              ["Intent Type", "Depth", "Multi-entity", "Alasan"],
              [
                  ["explain / followup",                   "2 hop", "—",  "Definisi resource cukup di depth 2"],
                  ["generate_yaml",                        "3 hop", "✓",  "Perlu field Container (depth 3) + Secret/ConfigMap"],
                  ["trace_relationship",                   "3 hop", "✓",  "Jembatan lintas-resource pada depth 2–3"],
                  ["planning",                             "3 hop", "✓",  "Multi-resource: hingga 2 related_concepts"],
              ],
              left=0.4, top=5.0, width=12.5,
              col_widths=[3.5, 1.5, 2, 5.5], font_size=12)


def slide_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "Dataset Evaluasi",
                   "97 Fixture Expert-Validated")

    # Left — table
    add_table(slide,
              ["Kategori", "n", "Deskripsi"],
              [
                  ["realworld",      "24", "Pertanyaan dari Stack Overflow (score > 5)"],
                  ["relationship",   "18", "Multi-hop relasi antar-resource"],
                  ["conceptual",     "15", "'Apa itu X?' — definitional"],
                  ["yaml_gen",       "15", "'Buat YAML untuk X' — schema compliance dinilai"],
                  ["followup",       "12", "Modifikasi / ekstensi konfigurasi sebelumnya"],
                  ["troubleshooting","5",  "Diagnosis error pod/workload"],
                  ["planning",       "5",  "Arsitektur multi-resource"],
                  ["command",        "3",  "Perintah kubectl"],
                  ["TOTAL",          "97", ""],
              ],
              left=0.4, top=1.75, width=7.5,
              col_widths=[3, 1, 5.5], font_size=12)

    # Right — validation info
    add_rect(slide, 8.2, 1.75, 4.9, 3.5, C_DARK_BLUE)
    txbox(slide, "Validasi Dataset", 8.4, 1.85, 4.5, 0.4,
          font_size=14, bold=True, color=C_ACCENT_TEAL)
    items_v = [
        "3 pakar DevOps/SRE",
        "Rata-rata realisme: 3,87/5,00",
        "Setiap fixture mencakup:",
        "  · ground truth answer",
        "  · relevant_nodes (KG paths)",
        "  · expected_path edges",
        "  · required_fields YAML",
        "Source: kubernetes.io docs",
        "  + Stack Overflow (realworld)",
    ]
    y = 2.35
    for item in items_v:
        txbox(slide, item, 8.4, y, 4.5, 0.35, font_size=12, color=C_TEXT_WHITE)
        y += 0.33

    # Fixture JSON preview
    txbox(slide, "Struktur Fixture (JSON):", 0.4, 5.45, 12.5, 0.3,
          font_size=12, bold=True, color=C_DARK_BLUE)
    code = '{ "id": "deployment_basic", "type": "yaml_gen", "question": "Buat YAML Deployment nginx …",\n  "ground_truth": { "relevant_nodes": […], "expected_path": ["Deployment -[HAS_PROPERTY]→ PodSpec"], … } }'
    txbox(slide, code, 0.4, 5.8, 12.5, 0.7,
          font_size=10, color=RGBColor(0x22, 0x44, 0x66))


def slide_eval_framework(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "Framework Evaluasi",
                   "3 Dimensi Custom: AnsQ · RetQ · ReaQ")

    # Formula box
    add_rect(slide, 0.4, 1.7, 12.5, 0.8, C_DARK_BLUE)
    txbox(slide,
          "Total Score  =  0,40 × AnsQ  +  0,35 × RetQ  +  0,25 × ReaQ",
          0.5, 1.8, 12.2, 0.6,
          font_size=20, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    dim_data = [
        ("AnsQ\nAnswer Quality", "40%", C_ACCENT_TEAL,
         ["Syntactic Validity (YAML safe_load)", "Schema Compliance (kubernetes-validate)",
          "Graph-Field Compliance (Neo4j)", "Faithfulness (node name overlap)",
          "Answer Relevance (cosine, text-embedding-3-small)"]),
        ("RetQ\nRetrieval Quality", "35%", RGBColor(0xE6, 0x7E, 0x22),
         ["Precision@k", "Recall@k", "F1@k",
          "Graph Coverage (% node relevan di-retrieve)", "NDCG@k", "Edge Coverage"]),
        ("ReaQ\nReasoning Quality", "25%", C_DARK_BLUE,
         ["Hop Accuracy", "Multi-Hop Success Rate",
          "Scope Accuracy (conditional)", "Grounding Score (canonical KG terms)",
          "Hallucination Rate (vocabulary-based)"]),
    ]
    x = 0.4
    for title, weight, col, metrics in dim_data:
        add_rect(slide, x, 2.65, 4.0, 0.75, col)
        txbox(slide, title, x + 0.1, 2.7, 2.5, 0.65,
              font_size=13, bold=True, color=C_TEXT_WHITE)
        txbox(slide, weight, x + 2.6, 2.7, 1.3, 0.65,
              font_size=24, bold=True, color=C_GOLD, align=PP_ALIGN.RIGHT)
        y = 3.5
        for m in metrics:
            txbox(slide, "· " + m, x + 0.05, y, 3.9, 0.35,
                  font_size=11, color=C_TEXT_DARK)
            y += 0.38
        x += 4.35


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "Hasil Perbandingan Sistem",
                   "GraphRAG unggul pada Retrieval Quality dan YAML Validity")

    # Main comparison table
    add_table(slide,
              ["Sistem", "AnsQ", "RetQ", "ReaQ", "Total Score"],
              [
                  ["GraphRAG (sistem ini)", "0,5798", "0,6487 ★", "0,8718", "0,6769 ★"],
                  ["Vector RAG",            "0,6047 ★", "0,4249",  "0,9019 ★", "0,6161"],
                  ["Vanilla LLM",           "0,5597",  "0,0241",  "0,6283",  "0,3894"],
              ],
              left=0.4, top=1.75, width=12.5,
              col_widths=[3.5, 2, 2, 2, 3], font_size=14)

    # Highlight boxes
    highlights = [
        ("YAML\nSyntactic Validity", "1,0000", C_ACCENT_TEAL, "100% YAML valid\nberkat validasi 3 lapis"),
        ("Multi-Hop\nSuccess Rate",  "1,0000", C_ACCENT_TEAL, "Traversal selalu ≥1 hop"),
        ("RetQ vs\nVector RAG",      "+0,22",  RGBColor(0xE6,0x7E,0x22), "Keunggulan nyata\ngraph traversal"),
        ("Graph\nCoverage",          "0,81",   C_DARK_BLUE,  "81% node relevan\nberhasil di-retrieve"),
    ]
    x = 0.4
    for title, val, col, note in highlights:
        add_rect(slide, x, 3.85, 2.9, 1.6, col)
        txbox(slide, title, x + 0.1, 3.9, 2.7, 0.5,
              font_size=11, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)
        txbox(slide, val, x + 0.1, 4.35, 2.7, 0.65,
              font_size=28, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
        txbox(slide, note, x + 0.1, 4.95, 2.7, 0.45,
              font_size=10, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)
        x += 3.15

    txbox(slide,
          "Kelemahan: realworld score 0,53 (terendah — batasan 3-hop) · hallucination rate 0,34",
          0.4, 5.65, 12.5, 0.4,
          font_size=12, italic=True, color=RGBColor(0x88, 0x44, 0x00))


def slide_iterations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "Iterasi Perbaikan",
                   "Riwayat Evaluasi v5 → v9 (versi produksi terkini: v9 = 0,6801)")

    add_table(slide,
              ["Versi", "Total Score", "Keterangan"],
              [
                  ["v5 (baseline thesis)", "0,6769", "Setelah expert validation — referensi utama TA"],
                  ["v6",                  "~0,66",  "Fix 1–5: speaker guard, retriever depth"],
                  ["v7",                  "~0,67",  "Fix 6–11: YAML rules, schema compliance, CronJob fix"],
                  ["v8 (invalid)",        "RetQ=0,02", "Neo4j mati saat evaluasi — data tidak valid"],
                  ["v9 ★",               "0,6801",  "Fix 12–14: reliability + OOD correction → melampaui baseline"],
              ],
              left=0.4, top=1.75, width=12.5,
              col_widths=[3, 2.5, 7], font_size=13)

    txbox(slide, "Fix Kunci (v9):", 0.4, 4.2, 12.5, 0.35,
          font_size=14, bold=True, color=C_DARK_BLUE)
    fixes = [
        ("Fix 12", "Tuple error strings untuk per-fixture retry detection di evaluate.py"),
        ("Fix 13", "Cegah false positive OOD rejection ketika Retrieved Data kosong — core K8s ALWAYS in-domain"),
        ("Fix 14", "Namespace exception untuk cluster-scoped resources (ClusterRoleBinding, PersistentVolume)"),
        ("Fix 16", "Multi-entity retrieval untuk trace_relationship + generate_yaml (extend dari planning-only)"),
    ]
    y = 4.65
    for code, desc in fixes:
        add_rect(slide, 0.4, y, 1.5, 0.38, C_ACCENT_TEAL)
        txbox(slide, code, 0.42, y + 0.04, 1.45, 0.3,
              font_size=12, bold=True, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)
        txbox(slide, desc, 2.05, y + 0.04, 10.8, 0.35,
              font_size=12, color=C_TEXT_DARK)
        y += 0.48


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "Kesimpulan",
                   "Tiga Rumusan Masalah Terjawab")

    rqs = [
        ("RQ1", C_ACCENT_TEAL,
         "Knowledge Graph berhasil dibangun dari swagger.json v1.30",
         "725 node · 18 jenis edge · 7 kategori relasi\nTaksonomi edge khusus merepresentasikan dependensi skema Kubernetes secara komprehensif."),
        ("RQ2", RGBColor(0xE6, 0x7E, 0x22),
         "Cascading Retrieval meningkatkan presisi secara signifikan",
         "RetQ = 0,65 (+0,22 di atas Vector RAG · +0,63 di atas Vanilla LLM)\nYAML Syntactic Validity = 1,0000 berkat validasi 3 lapis berbasis KG."),
        ("RQ3", C_DARK_BLUE,
         "GraphRAG meraih skor komposit tertinggi (0,6769 → v9: 0,6801)",
         "Keunggulan RetQ mengimbangi selisih kecil AnsQ dan ReaQ vs. Vector RAG.\nKeterbatasan: hallucination 0,34 · realworld score 0,53 (batasan 3-hop)."),
    ]
    y = 1.8
    for rq, col, title, body in rqs:
        add_rect(slide, 0.4, y, 1.0, 1.25, col)
        txbox(slide, rq, 0.4, y + 0.3, 1.0, 0.6,
              font_size=20, bold=True, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, 1.45, y, 11.4, 1.25, RGBColor(0xEE, 0xF2, 0xF8))
        txbox(slide, title, 1.55, y + 0.08, 11.1, 0.4,
              font_size=14, bold=True, color=col)
        txbox(slide, body, 1.55, y + 0.5, 11.1, 0.7,
              font_size=12, color=C_TEXT_DARK)
        y += 1.45

    txbox(slide, "Saran: adaptive traversal depth · constrained decoding · incremental KG update · Kubernetes Watch API integration",
          0.4, 6.05, 12.5, 0.45,
          font_size=12, italic=True, color=RGBColor(0x33, 0x55, 0x77))


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_DARK_BLUE)
    add_rect(slide, 0, 2.8, 13.33, 0.12, C_ACCENT_TEAL)
    add_rect(slide, 0, 5.3, 13.33, 0.12, C_ACCENT_TEAL)

    txbox(slide, "Terima Kasih", 0.5, 1.1, 12.3, 1.1,
          font_size=42, bold=True, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)
    txbox(slide, "Jihan Aurelia · 18222001 · jihanaurelia.jiji@gmail.com",
          0.5, 3.05, 12.3, 0.55,
          font_size=14, color=C_ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # Quick stats recap
    stats_recap = [
        ("725", "Node KG"),
        ("18", "Jenis Edge"),
        ("97", "Fixture Evaluasi"),
        ("0,6801", "Best Score (v9)"),
        ("1,0000", "YAML Validity"),
    ]
    x = 1.0
    for val, label in stats_recap:
        add_rect(slide, x, 3.9, 2.0, 1.2, RGBColor(0x0D, 0x1A, 0x35))
        txbox(slide, val, x + 0.05, 3.95, 1.9, 0.65,
              font_size=22, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
        txbox(slide, label, x + 0.05, 4.55, 1.9, 0.45,
              font_size=11, color=C_TEXT_WHITE, align=PP_ALIGN.CENTER)
        x += 2.25

    txbox(slide, "GraphRAG · Neo4j · LangGraph · GPT-4o-mini · Groq LLaMA · Streamlit",
          0.5, 5.5, 12.3, 0.45,
          font_size=12, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_problem(prs)
    slide_rq(prs)
    slide_architecture(prs)
    slide_kg(prs)
    slide_retrieval(prs)
    slide_dataset(prs)
    slide_eval_framework(prs)
    slide_results(prs)
    slide_iterations(prs)
    slide_conclusion(prs)
    slide_closing(prs)

    out = "PPT_Brief_TA_GraphRAG_Kubernetes.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
